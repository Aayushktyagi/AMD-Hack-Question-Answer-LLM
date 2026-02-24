#!/usr/bin/env python3
"""Generate a summary PowerPoint for the AAIPL project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────
BG_DARK        = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
BG_MEDIUM      = RGBColor(0x22, 0x22, 0x3A)
ACCENT_RED     = RGBColor(0xED, 0x1C, 0x24)   # AMD red
ACCENT_ORANGE  = RGBColor(0xFF, 0x7F, 0x32)
ACCENT_GREEN   = RGBColor(0x00, 0xC8, 0x53)
ACCENT_BLUE    = RGBColor(0x44, 0x8A, 0xFF)
ACCENT_PURPLE  = RGBColor(0xBB, 0x86, 0xFC)
TEXT_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT      = RGBColor(0xCC, 0xCC, 0xCC)
TEXT_DIM        = RGBColor(0x99, 0x99, 0x99)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── helpers ──────────────────────────────────────────────────────────
def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # round corners
    shape.adjustments[0] = 0.05
    return shape

def _add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def _set_text(tf, text, size=18, bold=False, color=TEXT_WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return p

def _add_bullet_frame(slide, left, top, width, height, items, size=16, color=TEXT_LIGHT, spacing=Pt(8), bold_prefix=True, icon="▸"):
    """Add a text box with bullet items.  Items can be plain strings or (bold_part, rest) tuples."""
    tb = _add_text_box(slide, left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
        if isinstance(item, tuple):
            r1 = p.add_run()
            r1.text = f"{icon} {item[0]}"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = TEXT_WHITE
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = f" {item[1]}"
            r2.font.size = Pt(size)
            r2.font.bold = False
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = f"{icon} {item}"
            r.font.size = Pt(size)
            r.font.bold = False
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb

def _add_card(slide, left, top, width, height, title, bullets, card_color=BG_MEDIUM, title_color=ACCENT_ORANGE):
    """Card = rounded rect + title + bullets."""
    _add_shape(slide, left, top, width, height, card_color)
    # title
    tb = _add_text_box(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.5))
    _set_text(tb.text_frame, title, size=18, bold=True, color=title_color)
    # bullets
    _add_bullet_frame(slide,
                      left + Inches(0.25), top + Inches(0.6),
                      width - Inches(0.5), height - Inches(0.75),
                      bullets, size=14, color=TEXT_LIGHT, icon="•")

# =====================================================================
# SLIDE 1 – Title
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
_set_slide_bg(slide, BG_DARK)

# Accent stripe
_add_shape(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(0.06), ACCENT_RED)

# Title
tb = _add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2))
_set_text(tb.text_frame, "AMD AI Premier League (AAIPL)", size=44, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
tb = _add_text_box(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.8))
_set_text(tb.text_frame, "Adversarial Agent Generation — Question & Answer Agents", size=24, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Tagline
tb = _add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6))
_set_text(tb.text_frame, "SFT + GRPO on Qwen & LLaMA  ·  AMD Instinct MI300X (192 GB HBM3)", size=18, bold=False, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Team
tb = _add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "Team AAIPL  |  Track 1", size=16, bold=False, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2 – Problem Statement
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Problem Statement", size=36, bold=True, color=ACCENT_RED)

_add_bullet_frame(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5), [
    ("Goal:", "Build a Q-Agent (generates hard MCQs) and an A-Agent (answers MCQs robustly) that compete adversarially."),
    ("Domain:", "Logical Reasoning — Syllogisms, Number/Alphanumeric Series, Blood Relations, Seating Arrangements."),
    ("Scoring:", "Q-Agent earns points by fooling the opponent's A-Agent; A-Agent earns points by correctly answering opponent's questions."),
    ("Constraint:", 'Questions must be valid & solvable — an "Oracle" (70B judge) verifies correctness.'),
], size=18, icon="➤")

# Two-column: Q-Agent vs A-Agent
_add_card(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0),
          "Q-Agent  (Offence)", [
              "Generate difficult but fair MCQs",
              "Exploit reasoning gaps of opponent",
              "4 choices (A–D), one correct answer",
              "Must pass Oracle validity check",
          ], title_color=ACCENT_ORANGE)

_add_card(slide, Inches(6.9), Inches(4.0), Inches(5.5), Inches(3.0),
          "A-Agent  (Defence)", [
              "Robustly answer adversarial MCQs",
              "Handle tricky distractors & traps",
              "Structured JSON output: {answer, explanation}",
              "Maximise accuracy under adversarial input",
          ], title_color=ACCENT_BLUE)

# =====================================================================
# SLIDE 3 – Topics & Sub-categories
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Topic Taxonomy & Sub-Categories", size=36, bold=True, color=ACCENT_RED)

topics_data = [
    ("Logical Reasoning / Syllogisms", ACCENT_ORANGE, [
        "both_neither_conclusion",
        "which_conclusion_follows",
        "which_does_not_follow",
    ]),
    ("Series & Patterns", ACCENT_BLUE, [
        "numeric_next_term",
        "alphanumeric_next_term",
        "missing_term",
        "odd_one_out",
    ]),
    ("Blood Relations & Family Tree", ACCENT_GREEN, [
        "simple_relation_2hop",
        "moderate_relation_3hop",
        "complex_relation_4hop",
        "extended_relation_5hop",
    ]),
    ("Puzzles / Seating Arrangements", ACCENT_PURPLE, [
        "linear_position_query",
        "linear_adjacent_query",
        "circular_position_query",
        "circular_adjacent_query",
        "circular_between_query",
    ]),
]

x_start = Inches(0.5)
card_w = Inches(2.95)
gap = Inches(0.15)
for idx, (title, color, subtypes) in enumerate(topics_data):
    left = x_start + idx * (card_w + gap)
    _add_shape(slide, left, Inches(1.4), card_w, Inches(5.5), BG_MEDIUM)
    # card title
    tb = _add_text_box(slide, left + Inches(0.15), Inches(1.55), card_w - Inches(0.3), Inches(0.7))
    _set_text(tb.text_frame, title, size=15, bold=True, color=color)
    # subtypes
    _add_bullet_frame(slide,
                      left + Inches(0.15), Inches(2.35),
                      card_w - Inches(0.3), Inches(4.3),
                      subtypes, size=13, color=TEXT_LIGHT, icon="•")

# =====================================================================
# SLIDE 4 – Dataset Generation Pipeline
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Dataset Generation Pipeline", size=36, bold=True, color=ACCENT_RED)

steps = [
    ("1", "Topic Definition", "Define 4 parent topics\n& 16 subtypes from\nlogical reasoning\ndomain", ACCENT_ORANGE),
    ("2", "Seed Generation", "Use LLaMA-3.3-70B\n(Oracle) to generate\ndiverse seed MCQs\nper subtype", ACCENT_BLUE),
    ("3", "Quality Filtering", "Oracle self-verification:\nsolve 5× at temp > 0;\nkeep if 4/5 agree\non correct answer", ACCENT_GREEN),
    ("4", "Format & Tag", "Tag each sample with\ntopic, subtype, difficulty;\nstructure as ChatML\nfor SFT training", ACCENT_PURPLE),
    ("5", "Final Dataset", "113,642 training\nsamples in\nqagent_chatml_train.json\n& aagent_chatml_val.json", ACCENT_RED),
]

box_w = Inches(2.2)
box_h = Inches(3.8)
start_x = Inches(0.4)
gap_x = Inches(0.25)
top_y = Inches(1.6)

for idx, (num, title, desc, color) in enumerate(steps):
    left = start_x + idx * (box_w + gap_x)
    # card bg
    _add_shape(slide, left, top_y, box_w, box_h, BG_MEDIUM)
    # number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.75), top_y + Inches(0.2), Inches(0.6), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    btf = badge.text_frame
    btf.word_wrap = False
    _set_text(btf, num, size=22, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # title
    tb = _add_text_box(slide, left + Inches(0.1), top_y + Inches(0.95), box_w - Inches(0.2), Inches(0.5))
    _set_text(tb.text_frame, title, size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    # desc
    tb = _add_text_box(slide, left + Inches(0.15), top_y + Inches(1.5), box_w - Inches(0.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    _set_text(tf, desc, size=13, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Arrow connectors (simple text arrows between boxes)
for idx in range(len(steps) - 1):
    left = start_x + (idx + 1) * (box_w + gap_x) - Inches(0.25)
    tb = _add_text_box(slide, left, top_y + Inches(1.5), Inches(0.3), Inches(0.5))
    _set_text(tb.text_frame, "→", size=28, bold=True, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Bottom note
tb = _add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "All data generated synthetically using LLaMA-3.3-70B-Instruct on AMD MI300X  •  No human annotation required",
          size=14, bold=False, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 5 – Answer Agent: SFT on Qwen
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "A-Agent Training: Supervised Fine-Tuning (SFT)", size=36, bold=True, color=ACCENT_RED)

# Left card – Model & Config
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5),
          "Model & Configuration", [
              ("Base Model:", "Qwen2.5-14B-Instruct"),
              ("Method:", "LoRA (rank 64, alpha 16)"),
              ("Precision:", "BF16 (full precision)"),
              ("Framework:", "Unsloth + HuggingFace TRL SFTTrainer"),
              ("Hardware:", "AMD MI300X (192 GB HBM3, ROCm)"),
              ("Training Data:", "113,642 ChatML samples (aagent_chatml_train.json)"),
              ("Format:", "system + user (topic/question/choices) → assistant (JSON answer)"),
              ("Epochs:", "3 full epochs, cosine LR schedule"),
          ], title_color=ACCENT_BLUE)

# Right card – Why SFT?
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.6),
          "Why SFT for A-Agent?", [
              "A-Agent needs reliable, structured output (JSON)",
              "SFT learns the exact format from curated data",
              "Consistent {answer, explanation} generation",
              "Robust against adversarial question phrasing",
          ], title_color=ACCENT_GREEN)

# Right card bottom – Output format
_add_card(slide, Inches(6.8), Inches(4.3), Inches(5.8), Inches(2.6),
          "Output Format (JSON)", [
              '{ "answer": "B",',
              '  "explanation": "Step-by-step reasoning..." }',
              "",
              "Robust parser handles malformed outputs",
              "6-strategy fallback chain in production",
          ], title_color=ACCENT_ORANGE)

# =====================================================================
# SLIDE 6 – Question Agent: GRPO Training
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Q-Agent Training: GRPO (Reinforcement Learning)", size=36, bold=True, color=ACCENT_RED)

# Left – GRPO overview
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5),
          "Group Relative Policy Optimization", [
              ("Base Model:", "Qwen2.5-14B-Instruct (LoRA rank 64)"),
              ("Framework:", "TRL GRPOTrainer + Unsloth + vLLM"),
              ("Oracle:", "LLaMA-3.3-70B-Instruct (4-bit, frozen)"),
              ("Generations/prompt:", "16 candidates per prompt"),
              ("Max Steps:", "1,500 steps"),
              ("GPU Memory:", "60% vLLM + 40% training"),
              ("Gradient Accum:", "1 (maximise throughput)"),
              ("Key Insight:", "No separate Critic needed — uses group-relative baseline"),
          ], title_color=ACCENT_ORANGE)

# Right top – Reward functions
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(2.8),
          "Reward Functions", [
              ("Format Reward:", "+1.0 for valid JSON, +1.0 for 4 choices, +1.0 for answer field"),
              ("Oracle Reward:", "+3.0 if Oracle verifies correct answer"),
              ("Wrong Penalty:", "−1.5 softened penalty (avoids collapse)"),
              ("Adversarial Bonus:", "+2.0 if Oracle is fooled (hard questions)"),
          ], title_color=ACCENT_GREEN)

# Right bottom – Training loop
_add_card(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.4),
          "GRPO Training Loop", [
              "1. Q-Agent generates 16 candidate MCQs per prompt",
              "2. Format reward: check JSON validity & structure",
              "3. Oracle (70B) attempts to solve each question",
              "4. Compute advantage relative to group baseline",
              "5. Policy gradient update (PPO-style, no critic)",
          ], title_color=ACCENT_PURPLE)

# =====================================================================
# SLIDE 7 – Architecture Overview (System Diagram)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "System Architecture", size=36, bold=True, color=ACCENT_RED)

# ── Q-Agent box ──
_add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.5), Inches(4.5), BG_MEDIUM)
tb = _add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.1), Inches(0.5))
_set_text(tb.text_frame, "Q-Agent (Offence)", size=20, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
_add_bullet_frame(slide, Inches(0.7), Inches(2.2), Inches(3.1), Inches(3.5), [
    "Qwen2.5-14B + LoRA",
    "GRPO-trained policy",
    "Subtype-aware prompts",
    "16 subtypes across 4 topics",
    "Robust JSON parser (5 strategies)",
    "Overgenerate → Filter → Rank",
], size=13, color=TEXT_LIGHT, icon="•")

# ── Oracle box ──
_add_shape(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(3.5), BG_MEDIUM)
tb = _add_text_box(slide, Inches(4.9), Inches(2.1), Inches(3.4), Inches(0.5))
_set_text(tb.text_frame, "Oracle / Judge (70B)", size=20, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
_add_bullet_frame(slide, Inches(4.9), Inches(2.7), Inches(3.4), Inches(2.5), [
    "LLaMA-3.3-70B-Instruct",
    "4-bit quantized (frozen)",
    "Validates Q-Agent output",
    "Provides reward signal",
    "Batched inference for speed",
], size=13, color=TEXT_LIGHT, icon="•")

# ── A-Agent box ──
_add_shape(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(4.5), BG_MEDIUM)
tb = _add_text_box(slide, Inches(9.5), Inches(1.6), Inches(3.1), Inches(0.5))
_set_text(tb.text_frame, "A-Agent (Defence)", size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
_add_bullet_frame(slide, Inches(9.5), Inches(2.2), Inches(3.1), Inches(3.5), [
    "Qwen2.5-14B + LoRA",
    "SFT-trained (113K samples)",
    "Structured JSON output",
    "Robust parser (6 strategies)",
    "Verification-first prompting",
    "Handles adversarial inputs",
], size=13, color=TEXT_LIGHT, icon="•")

# Arrows
for (x, label) in [(Inches(4.1), "generates →"), (Inches(8.6), "← answers")]:
    tb = _add_text_box(slide, x, Inches(3.4), Inches(1.5), Inches(0.5))
    _set_text(tb.text_frame, label, size=14, bold=True, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Hardware footer
_add_shape(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.8), RGBColor(0x2A, 0x2A, 0x44))
tb = _add_text_box(slide, Inches(1.7), Inches(6.35), Inches(9.9), Inches(0.7))
_set_text(tb.text_frame, "Hardware:  AMD Instinct MI300X  ·  192 GB HBM3  ·  5.3 TB/s bandwidth  ·  ROCm + vLLM",
          size=16, bold=False, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 8 – Robust Output Parsing
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Robust Output Parsing", size=36, bold=True, color=ACCENT_RED)

# Q-Agent parser
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.2),
          "Q-Agent Parser  (5 Strategies)", [
              ("1. Direct JSON:", 'json.loads() on raw output'),
              ("2. Code-block extraction:", "Strip ```json fences, parse inner text"),
              ("3. Brace extraction:", "Find outermost { … } and parse"),
              ("4. Field-per-line:", 'Regex: key"value" per line → dict'),
              ("5. Partial key extraction:", 'Regex "topic", "question", "0"-"3" fields'),
          ], title_color=ACCENT_ORANGE)

# A-Agent parser
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
          "A-Agent Parser  (6 Strategies)", [
              ("1. Direct JSON:", 'json.loads() on raw output'),
              ("2. Null-line strip:", 'Remove \\x00 / NUL bytes then parse'),
              ("3. Code-block extraction:", "Strip ```json fences"),
              ("4. Brace extraction:", "Find outermost { … }"),
              ("5. Key-value regex:", 'Match "answer":"X" pattern'),
              ("6. First-letter fallback:", "First A/B/C/D character in output"),
          ], title_color=ACCENT_BLUE)

# Bottom note
tb = _add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "3-tier pipeline:  Direct JSON  →  Robust Parser  →  LLM Re-prompt Fallback",
          size=16, bold=False, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 9 – Evaluation & Metrics
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Evaluation & Key Metrics", size=36, bold=True, color=ACCENT_RED)

# Metrics table simulation
metrics = [
    ("Format Compliance", "3.0 / 3.0", "Valid JSON with correct structure", ACCENT_GREEN),
    ("Length Compliance", "3.0 / 3.0", "Within token limits", ACCENT_GREEN),
    ("Oracle + Adversarial Mean", "1.28", "Average oracle & adversarial score", ACCENT_ORANGE),
    ("High Oracle %", "0%", "% questions Oracle couldn't solve (target: ↑)", ACCENT_RED),
    ("Oracle Fooled %", "26%", "% questions that tricked opponent A-Agent", ACCENT_ORANGE),
]

row_h = Inches(0.85)
y_start = Inches(1.5)
for idx, (metric, value, desc, color) in enumerate(metrics):
    y = y_start + idx * (row_h + Inches(0.1))
    _add_shape(slide, Inches(0.8), y, Inches(11.5), row_h, BG_MEDIUM)
    # metric name
    tb = _add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(3.5), Inches(0.5))
    _set_text(tb.text_frame, metric, size=17, bold=True, color=color)
    # value
    tb = _add_text_box(slide, Inches(4.8), y + Inches(0.15), Inches(1.8), Inches(0.5))
    _set_text(tb.text_frame, value, size=20, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    # description
    tb = _add_text_box(slide, Inches(6.8), y + Inches(0.15), Inches(5.0), Inches(0.5))
    _set_text(tb.text_frame, desc, size=14, bold=False, color=TEXT_LIGHT)

# Footer
tb = _add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5))
_set_text(tb.text_frame, "Metrics from initial GRPO checkpoint (step 51)  •  Training continued to 1,500 steps with improved rewards",
          size=14, bold=False, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 10 – Improvements Applied
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Key Improvements & Iterations", size=36, bold=True, color=ACCENT_RED)

# Left column
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5),
          "Training Improvements", [
              ("Batched Oracle Inference:", "10× faster reward computation"),
              ("Gradient Accum 4→1:", "More frequent updates, better signal"),
              ("Generations 8→16:", "Better advantage estimation in GRPO"),
              ("Max Steps 500→1500:", "Longer training for convergence"),
              ("Softened Penalty −3→−1.5:", "Avoids reward collapse, encourages exploration"),
              ("Resume from Checkpoint:", "Continuous training without data loss"),
          ], title_color=ACCENT_ORANGE)

# Right column
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.5),
          "Inference & Evaluation Improvements", [
              ("Subtype-Specific Prompts:", "Match exact SFT training format in eval"),
              ("Robust Parsers:", "5-strategy (Q) + 6-strategy (A) fallback chains"),
              ("In-place Answer Parsing:", "_parse_and_normalize() inside answer_batches()"),
              ("HF generate() Fix:", "Replaced broken fast_generate() with HF kwargs"),
              ("Left-Pad Batched Gen:", "Proper padding for batch inference"),
              ("Standalone Validation:", "run_model.py with confusion matrices"),
          ], title_color=ACCENT_BLUE)

# =====================================================================
# SLIDE 11 – Hardware & Stack
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Hardware & Software Stack", size=36, bold=True, color=ACCENT_RED)

# Hardware card
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(3.0),
          "AMD Instinct MI300X", [
              ("VRAM:", "192 GB HBM3"),
              ("Bandwidth:", "5.3 TB/s"),
              ("Architecture:", "CDNA 3 chiplet design"),
              ("Advantage:", "Full 70B model in BF16 on single GPU — no tensor parallelism needed"),
          ], title_color=ACCENT_RED)

# Software card
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3.0),
          "Software Stack", [
              ("Runtime:", "ROCm + vLLM (PagedAttention)"),
              ("Training:", "Unsloth + HuggingFace TRL"),
              ("Models:", "Qwen2.5-14B, LLaMA-3.3-70B"),
              ("Precision:", "BF16 (training) / 4-bit (Oracle inference)"),
          ], title_color=ACCENT_BLUE)

# Memory allocation diagram (text-based)
_add_card(slide, Inches(0.5), Inches(4.7), Inches(12.1), Inches(2.3),
          "GPU Memory Allocation During GRPO Training", [
              ("vLLM Engine (60%):", "~115 GB — Hosts Q-Agent policy for fast generation of 16 candidates/prompt"),
              ("Training (40%):", "~77 GB — LoRA gradients, optimizer states, reference model, reward computation"),
              ("Oracle (separate):", "LLaMA-3.3-70B in 4-bit (~35 GB) loaded for batched reward evaluation"),
          ], title_color=ACCENT_GREEN)

# =====================================================================
# SLIDE 12 – Summary & Next Steps
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
_set_slide_bg(slide, BG_DARK)

tb = _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
_set_text(tb.text_frame, "Summary & Next Steps", size=36, bold=True, color=ACCENT_RED)

# Summary
_add_card(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.2),
          "What We Built", [
              "Synthetic dataset: 113K MCQs across 4 topics, 16 subtypes",
              "A-Agent: SFT-trained Qwen2.5-14B (LoRA) — reliable JSON answerer",
              "Q-Agent: GRPO-trained Qwen2.5-14B — adversarial question generator",
              "Oracle: LLaMA-3.3-70B for validation & reward computation",
              "Robust 5+6 strategy parsers for production reliability",
              "Full pipeline on single AMD MI300X (192 GB)",
          ], title_color=ACCENT_GREEN)

# Next steps
_add_card(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2),
          "Next Steps", [
              "Continue GRPO to 1,500 steps — target higher Oracle-fooled %",
              "Hard subtype focus (odd_one_out, 5-hop relations, circular_between)",
              "Distractor quality: Overgenerate-and-Rank pipeline",
              "A-Agent: Test-Time Compute scaling (Best-of-N + Verifier)",
              "Self-play co-evolution: Q-Agent vs A-Agent arena",
              "Ablation studies on reward function weights",
          ], title_color=ACCENT_PURPLE)

# Footer
_add_shape(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.06), ACCENT_RED)
tb = _add_text_box(slide, Inches(1), Inches(7.1), Inches(11), Inches(0.4))
_set_text(tb.text_frame, "AMD AI Premier League  ·  Track 1  ·  Adversarial Agent Generation",
          size=12, bold=False, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────
out_path = "AAIPL_Summary.pptx"
prs.save(out_path)
print(f"✅ Saved presentation → {out_path}")
print(f"   {len(prs.slides)} slides")
